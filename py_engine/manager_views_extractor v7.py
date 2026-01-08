"""
EXHAUSTIVE MANAGER INSIGHT EXTRACTOR (DOCX / PDF / PPTX -> DOCX)

Goal:
- Given an input investment memo / deck (DOCX, PDF, PPTX) and a manager name,
  produce an exhaustive, well-organized bullet profile about that manager
  (similar depth to what was done in chat), and save it to a Word (.docx).

What this script does (simple + high-recall):
1) Extracts text from DOCX/PDF/PPTX.
2) Chunks the text.
3) Runs an LLM "map" pass over each chunk to extract ALL manager-related info.
4) Runs a "reduce" pass to merge, deduplicate, and organize into a final bullet profile.
5) Writes the final output to a .docx file.

Requirements:
- Python 3.10+
- pip install pymupdf python-docx python-pptx openai

Run:
  export OPENAI_API_KEY="your_key"
  python manager_extract.py --input "/path/to/file.pdf" --manager "Asana" --output "Asana_Manager_Profile.docx"

Notes:
- No citations/page refs by design (per request).
- If you later want citations, we can add them back easily.
"""

import os
import re
import math
import argparse
import json
import datetime
from dataclasses import dataclass
from typing import List, Optional, Tuple

# Input parsers
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

# Output writer
from docx import Document as DocxDocument
from docx.shared import Pt

# LLM
from openai import OpenAI


# ----------------------------
# Output configuration
# ----------------------------
# Set this to a directory where outputs should be written.
# If the directory does not exist or is empty/None, the script will fall back to the input file directory.
DEFAULT_OUTPUT_DIR = r"C:\Users\prash\Downloads"


# ----------------------------
# Utilities
# ----------------------------



def parse_document_source(raw: str) -> str:
    """Parse --document_source. Accepts short codes (M/T/3) or full labels.

    Canonical outputs:
      - Manager
      - Townsend
      - 3rd party
    """
    if raw is None:
        raise argparse.ArgumentTypeError("Missing document source.")
    s = str(raw).strip()
    if not s:
        raise argparse.ArgumentTypeError("Document source cannot be empty.")
    key = s.casefold().replace("_", " ").strip()

    mapping = {
        "m": "Manager",
        "manager": "Manager",
        "t": "Townsend",
        "townsend": "Townsend",
        "3": "3rd party",
        "3rd": "3rd party",
        "3rd party": "3rd party",
        "third": "3rd party",
        "third party": "3rd party",
        "thirdparty": "3rd party",
        "3rdparty": "3rd party",
    }
    if key in mapping:
        return mapping[key]
    raise argparse.ArgumentTypeError('Invalid --document_source. Use M/Manager, T/Townsend, or 3/"3rd party".')
def normalize_text(text: str) -> str:
    if not text:
        return ""
    # Normalize whitespace while keeping paragraph boundaries
    text = text.replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(text: str, max_chars: int = 12000, overlap_chars: int = 1200) -> List[str]:
    """
    Chunk text by character count with overlap to preserve cross-boundary context.
    """
    text = normalize_text(text)
    if not text:
        return []

    chunks = []
    start = 0
    n = len(text)
    while start < n:
        end = min(n, start + max_chars)
        chunk = text[start:end]
        chunks.append(chunk)
        if end == n:
            break
        start = max(0, end - overlap_chars)
    return chunks


# ----------------------------
# Extract text from formats
# ----------------------------

def extract_from_pdf(path: str) -> str:
    doc = fitz.open(path)
    parts = []
    for i in range(len(doc)):
        page = doc.load_page(i)
        t = page.get_text("text")
        t = normalize_text(t)
        if t:
            parts.append(t)
    return "\n\n".join(parts)


def extract_from_docx(path: str) -> str:
    d = Document(path)
    parts = []

    # Paragraph text
    for p in d.paragraphs:
        txt = p.text.strip()
        if txt:
            parts.append(txt)

    # Tables (simple flatten)
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    return normalize_text("\n\n".join(parts))


def extract_from_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = shape.text.strip()
                if txt:
                    slide_text.append(txt)
        if slide_text:
            parts.append("\n".join(slide_text))
    return normalize_text("\n\n".join(parts))


def extract_text_by_extension(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_from_pdf(path)
    if ext == ".docx":
        return extract_from_docx(path)
    if ext == ".pptx":
        return extract_from_pptx(path)
    raise ValueError(f"Unsupported input type: {ext}. Use .pdf, .docx, or .pptx")


# ----------------------------
# LLM prompts (exhaustive extraction)
# ----------------------------

SYSTEM_PROMPT = """You are an investment due diligence analyst.
Your job: extract EVERYTHING about the named manager from the provided document text.

Rules:
- Be exhaustive and high-recall: include every relevant fact, claim, metric, list, name, process detail, governance detail, risk, mitigation, and nuance.
- Do NOT invent anything. If unsure, omit.
- Output must be structured bullets, organized into clear sections.
- No citations and no page references.
- If the document has conflicting numbers/statements, keep BOTH and flag as "Inconsistency noted".
- Prefer concrete details over generic wording.
- If the text describes the fund more broadly, still capture manager-related implications (team, platform, capabilities, behavior, policies, etc.).
- Write in the style of a comprehensive manager profile.

Use these section headers (add more if needed):
1) Identity & Platform
2) Leadership, Ownership & Governance
3) Team, Roles & Incentives
4) Strategy & Differentiation
5) Investment & Asset Management Process
6) Products / Vehicles Managed
7) Track Record & Performance Notes
8) Operations, Controls & Service Providers
9) ESG / Sustainability / DEI (if present)
10) Risks / Concerns / Open Questions
11) Notable Claims & "Manager Says" Statements
"""

MAP_PROMPT_TEMPLATE = """Manager name: {manager}

TASK:
From the following document excerpt, extract ALL information that pertains to the manager {manager} (or its platform, affiliates, team, governance, capabilities, operations, track record, ESG, risks, etc.).

OUTPUT FORMAT:
- Use section headers (you can repeat headers across chunks).
- Under each header, use bullet points.
- Be exhaustive.
- If something is ambiguous but likely manager-related, include it as a bullet prefixed with "Possible:".

DOCUMENT EXCERPT:
{chunk}
"""

REDUCE_PROMPT_TEMPLATE = """Manager name: {manager}

You will be given multiple extracted bullet lists (from different excerpts of the same document).
Merge them into ONE master "Manager Profile" that is:
- exhaustive (high recall)
- well organized into sections
- deduplicated (no repeated bullets)
- internally consistent where possible, but if the source has conflicts, keep both and mark "Inconsistency noted"
- written in crisp, investment-committee-ready bullets

Do NOT include citations.

EXTRACTED BULLETS:
{all_maps}
"""


# ----------------------------
# LLM runner
# ----------------------------

@dataclass
class LLMConfig:
    model: str = "gpt-5.2"          # change if needed
    temperature: float = 0.1
    max_output_tokens: Optional[int] = None  # let API decide unless you want to cap


def call_llm(client: OpenAI, system: str, user: str, cfg: LLMConfig) -> str:
    """
    Uses OpenAI Responses API. Adjust if your environment uses a different client.
    """
    kwargs = {}
    if cfg.max_output_tokens is not None:
        kwargs["max_output_tokens"] = cfg.max_output_tokens

    resp = client.responses.create(
        model=cfg.model,
        temperature=cfg.temperature,
        input=[
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        **kwargs
    )
    return resp.output_text.strip()


def extract_manager_profile(text: str, manager: str, cfg: LLMConfig) -> str:
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("Missing OPENAI_API_KEY environment variable.")

    client = OpenAI(api_key=api_key)

    chunks = chunk_text(text, max_chars=6000, overlap_chars=2000)
    if not chunks:
        raise RuntimeError("No extractable text found in the input document.")

    # MAP: extract bullets from each chunk
    map_outputs = []
    for i, ch in enumerate(chunks, start=1):
        user_prompt = MAP_PROMPT_TEMPLATE.format(manager=manager, chunk=ch)
        out = call_llm(client, SYSTEM_PROMPT, user_prompt, cfg)
        map_outputs.append(out)

    # REDUCE: merge and organize
    all_maps = "\n\n---\n\n".join(map_outputs)
    reduce_prompt = REDUCE_PROMPT_TEMPLATE.format(manager=manager, all_maps=all_maps)
    final_profile = call_llm(client, SYSTEM_PROMPT, reduce_prompt, cfg)

    return final_profile


# ----------------------------
# Write output to DOCX

def write_profile_to_docx(manager: str, profile_text: str, out_path: str) -> None:
    doc = DocxDocument()

    # Styling
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)

    doc.add_heading(f"{manager} - Extracted Manager Profile", level=1)

    # Parse: treat lines starting with "-", "•", or "*" as bullets; others as headings/paragraphs
    for raw_line in profile_text.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        # Markdown-ish headings
        if re.match(r"^\d+\)\s+", line):
            # "1) Identity & Platform"
            doc.add_heading(line, level=2)
            continue
        if line.endswith(":") and len(line) <= 80:
            # "Identity & Platform:"
            doc.add_heading(line[:-1], level=2)
            continue

        if line.startswith(("-", "•", "*")):
            bullet = line.lstrip("-•* ").strip()
            doc.add_paragraph(bullet, style="List Bullet")
        else:
            # If LLM emits a header without numbering
            if re.match(r"^(Identity|Leadership|Team|Strategy|Investment|Products|Track Record|Operations|ESG|Risks|Notable)\b", line, re.I):
                doc.add_heading(line, level=2)
            else:
                doc.add_paragraph(line)

    doc.save(out_path)



# ----------------------------
# Write RAW output to JSON (no MCAF structuring)
# ----------------------------

def write_profile_to_raw_json(
    manager: str,
    manager_id: int,
    document_id: int,
    document_date: str,
    document_source: str,
    extraction_datetime: str,
    input_path: str,
    output_docx_path: str,
    profile_text: str,
    out_json_path: str,
) -> None:
    """
    Writes a RAW JSON artifact that is intentionally minimally structured.
    It is derived from the same in-memory `profile_text` used to generate the DOCX,
    and is intended to preserve *all* extracted facts without forcing an MCAF schema.
    """
    now = extraction_datetime
    payload = {
        "manager_id": manager_id,
        "manager_name": manager,
        "document_metadata": {
            "document_id": document_id,
            "document_date": document_date,
            "document_source": document_source,
            "extraction_datetime": extraction_datetime,
        },
        "extraction_metadata": {
            "input_file": os.path.abspath(input_path),
            "generated_from": os.path.basename(__file__),
            "output_docx": os.path.abspath(output_docx_path),
            "generated_at": now,
        },
        # Preserve the raw text exactly as produced by the LLM reduce pass
        "profile_text": profile_text,
        # Convenience: non-empty lines (still unstructured)
        "profile_lines": [ln.strip() for ln in profile_text.splitlines() if ln.strip()],
    }
    with open(out_json_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, indent=2, ensure_ascii=False)


# ----------------------------
# Main
# ----------------------------

def _build_output_paths(input_path: str, manager_id: int, document_id: int) -> tuple[str, str]:
    """
    Output files are created in the configured output directory (DEFAULT_OUTPUT_DIR) if it exists,
    otherwise in the SAME directory as the input file:
      <manager_id>_raw_<document_id>.docx
      <manager_id>_raw_<document_id>.json
    """
    preferred_dir = DEFAULT_OUTPUT_DIR.strip() if isinstance(DEFAULT_OUTPUT_DIR, str) else ""
    if preferred_dir and os.path.isdir(preferred_dir):
        out_dir = preferred_dir
    else:
        out_dir = os.path.dirname(os.path.abspath(input_path))

    base = f"{manager_id}_raw_{document_id}"
    out_docx = os.path.join(out_dir, base + ".docx")
    out_json = os.path.join(out_dir, base + ".json")
    return out_docx, out_json


def main():
    parser = argparse.ArgumentParser(description="LLM-assisted exhaustive manager extraction -> RAW DOCX + RAW JSON")
    parser.add_argument("--input", required=True, help="Path to input file (.pdf, .docx, .pptx)")
    parser.add_argument("--manager", required=True, help="Manager name to extract (e.g., Asana)")
    parser.add_argument("--manager_id", type=int, required=True, help="Manager ID (int). Used for output filenames.")
    parser.add_argument("--document_id", type=int, required=True, help="Document ID (int). Used for output filenames.")
    parser.add_argument("--document_date", required=True, help="Document date in YYYY-MM-DD format. Used for audit/tracing.")
    parser.add_argument("--document_source", required=True, type=parse_document_source, help='Document source: M/Manager, T/Townsend, or 3/"3rd party".')
    parser.add_argument("--model", default="gpt-5.2", help="OpenAI model name (default: gpt-5.2)")
    parser.add_argument("--temperature", type=float, default=0.1, help="LLM temperature (default: 0.1)")
    args = parser.parse_args()

    # Validate document_date (YYYY-MM-DD)
    try:
        _ = datetime.date.fromisoformat(args.document_date)
    except Exception:
        raise ValueError('Invalid --document_date. Expected format: YYYY-MM-DD (e.g., 2025-02-15)')


    # Canonical document source (validated by argparse)
    document_source = args.document_source

    # Extract text from source
    text = extract_text_by_extension(args.input)

    # LLM extraction (unchanged)
    cfg = LLMConfig(model=args.model, temperature=args.temperature)
    profile = extract_manager_profile(text=text, manager=args.manager, cfg=cfg)

    # Prepend audit header lines (document date and extraction date) to the profile text.
    # NOTE: The DOCX writer adds its own title heading before rendering this text; we do not modify DOCX logic.
    extraction_ts = datetime.datetime.now().isoformat(timespec="seconds")
    audit_header = f"Document Date: {args.document_date}\nExtraction Date: {extraction_ts}\nDocument Source: {document_source}\n\n"
    profile = audit_header + profile


    # Output paths in same directory as input (not CLI-dependent)
    out_docx, out_json = _build_output_paths(args.input, args.manager_id, args.document_id)

    # DOCX generation (DO NOT MODIFY)
    write_profile_to_docx(args.manager, profile, out_docx)
    print(f"Done. Wrote: {out_docx}")

    # RAW JSON generation (no MCAF structuring)
    write_profile_to_raw_json(
        manager=args.manager,
        manager_id=args.manager_id,
        document_id=args.document_id,
        document_date=args.document_date,
        document_source=document_source,
        extraction_datetime=extraction_ts,
        input_path=args.input,
        output_docx_path=out_docx,
        profile_text=profile,
        out_json_path=out_json,
    )
    print(f"Done. Wrote JSON: {out_json}")


if __name__ == "__main__":
    main()
